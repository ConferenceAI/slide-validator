import io

import zipfile
import io
import json
import base64
import xml.etree.ElementTree as ET

import re
import fitz
from pikepdf import Pdf
from PyPDF2 import PdfReader

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def analyze_file(content: bytes, filename: str, file_format) -> dict:
    analysis = {
        'slide_count': 0,
        'image_count': 0,
        'audio_count': 0,
        'video_count': 0,
        'bullet_count': 0,
        'fonts_used': [],
    }

    if file_format == "application/pdf":
        print("Analyzing PDF...")
        analyze_pdf(content, analysis)
    elif file_format == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        print("Analyzing PPTX...")
        analyze_pptx(content, analysis)
    elif file_format == "application/x-iwork-keynote-sffkey" or (
            file_format == "application/zip" and filename.endswith('.key')):
        print("Analyzing Keynote...")
        analyze_keynote(content, analysis)
    elif file_format == "application/json" and filename.endswith('.canva'):
        print("Analyzing Canva...")
        analyze_canva(content, analysis)
    elif file_format == "application/octet-stream" and filename.endswith(
            '.fig'):
        print("Analyzing Figma...")
        analyze_figma(content, analysis)
    elif file_format == "text/markdown" and filename.lower().endswith('.md'):
        print("Analyzing Markdown...")
        analyze_readme_md(content, analysis)

    return analysis


def analyze_pdf(content: bytes, analysis: dict):
    pdf_stream = io.BytesIO(content)
    pdf = PdfReader(pdf_stream)
    analysis['slide_count'] = len(pdf.pages)

    # Font analysis
    fonts_used = set()
    for page in pdf.pages:
        if '/Resources' in page and '/Font' in page['/Resources']:
            for font in page['/Resources']['/Font']:
                if '/BaseFont' in page['/Resources']['/Font'][font]:
                    fonts_used.add(
                        page['/Resources']['/Font'][font]['/BaseFont'])
    analysis['fonts_used'] = list(fonts_used)

    # Image, audio, and video analysis using pikepdf
    pdf_pike = Pdf.open(pdf_stream)
    image_count = 0
    audio_count = 0
    video_count = 0

    for page in pdf_pike.pages:
        for name, obj in page.get("/Resources", {}).get("/XObject",
                                                        {}).items():
            if obj.get("/Subtype") == "/Image":
                image_count += 1
            elif obj.get("/Subtype") == "/Sound":
                audio_count += 1
            elif obj.get("/Subtype") == "/Movie":
                video_count += 1

    analysis['image_count'] = image_count
    analysis['audio_count'] = audio_count
    analysis['video_count'] = video_count

    # Bullet point analysis
    bullet_count = 0
    doc = fitz.open(stream=content, filetype="pdf")

    bullet_patterns = [
        r'^\s*[\u2022\u2023\u25E6\u2043\u2219]\s',  # Unicode bullet characters
        r'^\s*[-•*]\s',  # Hyphen, bullet, asterisk
        r'^\s*\d+\.?\s',  # Numbered list (e.g., "1. " or "1 ")
        r'^\s*[a-zA-Z]\.?\s',  # Lettered list (e.g., "a. " or "a ")
        r'^\s*\(\d+\)\s',  # Parenthesized numbers (e.g., "(1) ")
        r'^\s*\([a-zA-Z]\)\s',  # Parenthesized letters (e.g., "(a) ")
    ]

    def is_bullet_point(text, x, page_width):
        # Check if the text matches any bullet pattern
        if any(re.match(pattern, text) for pattern in bullet_patterns):
            return True

        # Check if the text is positioned like a bullet point (near the left margin)
        if x < page_width * 0.1 and len(
                text.strip()
        ) < 5:  # Assuming bullets are within 10% of page width from the left and short
            return True

        return False

    for page in doc:
        blocks = page.get_text("dict")["blocks"]
        page_width = page.rect.width

        for block in blocks:
            if block.get("type") == 0:  # Type 0 is text
                previous_bullet_y = None
                for line in block.get("lines", []):
                    line_text = "".join(span["text"]
                                        for span in line.get("spans", []))
                    first_span = line.get("spans", [{}])[0]
                    x0 = first_span.get("origin", (0, 0))[0]
                    y0 = first_span.get("origin", (0, 0))[1]

                    if is_bullet_point(line_text, x0, page_width):
                        # Check if this bullet is significantly below the previous one
                        if previous_bullet_y is None or (
                                y0 - previous_bullet_y
                        ) > 5:  # Adjust the 5 as needed
                            bullet_count += 1
                            previous_bullet_y = y0

    analysis['bullet_count'] = bullet_count

    doc.close()
    return analysis


def analyze_pptx(content: bytes, analysis: dict):
    prs = Presentation(io.BytesIO(content))

    analysis['slide_count'] = len(prs.slides)
    analysis['image_count'] = 0
    analysis['audio_count'] = 0
    analysis['video_count'] = 0
    analysis['bullet_count'] = 0
    fonts_used = set()

    def is_bullet_point(paragraph):
        # Check if the paragraph has bullet characteristics
        if paragraph.level > 0:
            return True
        if hasattr(paragraph, 'bullet') and paragraph.bullet is not None:
            return True
        # Check if the paragraph starts with a bullet-like character
        text = paragraph.text.strip()
        if any(
                text.startswith(char) for char in ('•', '-', '*', '✓', '☐',
                                                   '○', '►', '➢', '➣', '➤')):
            return True
        # Check for numbered list (extended to handle multidigit and different separators)
        if any(text[:i].isdigit() and text[i:].startswith(sep)
               for i in range(1, 4) for sep in ('. ', ') ', ' ')):
            return True
        # Check for lettered list (e.g., "a.", "b)", "c ")
        if len(text) > 1 and text[0].isalpha() and text[1] in ('.', ')', ' '):
            return True
        # Check for parenthesized numbers
        if re.match(r'^\s*\(\d+\)\s', text):
            return True
        # Check for parenthesized letters
        if re.match(r'^\s*\([a-zA-Z]\)\s', text):
            return True
        return False

    def count_bullets_in_shape(shape):
        bullet_count = 0
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                if is_bullet_point(paragraph):
                    bullet_count += 1
        elif hasattr(shape, 'table'):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if is_bullet_point(paragraph):
                            bullet_count += 1
        return bullet_count

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                analysis['image_count'] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                if hasattr(shape, 'video'):
                    analysis['video_count'] += 1
                elif hasattr(shape, 'audio'):
                    analysis['audio_count'] += 1

            analysis['bullet_count'] += count_bullets_in_shape(shape)

            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts_used.add(run.font.name)

    analysis['fonts_used'] = list(fonts_used)
    return analysis


def analyze_keynote(content: bytes, analysis: dict):
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        # Parse index.apxl file for slide count and bullet points
        with zf.open('index.apxl') as f:
            tree = ET.parse(f)
            root = tree.getroot()

            # Count slides
            slides = root.findall(
                ".//sf:slide",
                namespaces={'sf': 'http://developer.apple.com/namespaces/sf'})
            analysis['slide_count'] = len(slides)

            # Count bullet points
            bullet_points = root.findall(
                ".//sf:bullet-characters",
                namespaces={'sf': 'http://developer.apple.com/namespaces/sf'})
            analysis['bullet_count'] = len(bullet_points)

        # Count images, audio, and video
        analysis['image_count'] = sum(
            1 for file in zf.namelist()
            if file.startswith('Data/') and file.lower().endswith((
                '.png', '.jpg', '.jpeg', '.gif')))
        analysis['audio_count'] = sum(
            1 for file in zf.namelist()
            if file.startswith('Data/') and file.lower().endswith(('.mp3',
                                                                   '.wav',
                                                                   '.m4a')))
        analysis['video_count'] = sum(
            1 for file in zf.namelist()
            if file.startswith('Data/') and file.lower().endswith(('.mp4',
                                                                   '.mov',
                                                                   '.m4v')))

        # Extract fonts used
        fonts_used = set()
        for file in zf.namelist():
            if file.endswith('.apxl'):
                with zf.open(file) as f:
                    content = f.read().decode('utf-8')
                    font_matches = re.findall(r'font-family="([^"]+)"',
                                              content)
                    fonts_used.update(font_matches)

        analysis['fonts_used'] = list(fonts_used)


def analyze_canva(content: bytes, analysis: dict):
    canva_data = json.loads(content)

    analysis['slide_count'] = len(canva_data.get('pages', []))
    analysis['image_count'] = 0
    analysis['bullet_count'] = 0
    analysis['fonts_used'] = set()

    def process_element(element):
        if element.get('type') == 'IMAGE':
            analysis['image_count'] += 1
        elif element.get('type') == 'TEXT':
            text = element.get('text', '')
            if text.strip().startswith(('•', '-', '*')):
                analysis['bullet_count'] += 1
            if 'font' in element:
                analysis['fonts_used'].add(element['font'])

    for page in canva_data.get('pages', []):
        for element in page.get('elements', []):
            process_element(element)

    analysis['fonts_used'] = list(analysis['fonts_used'])


def analyze_figma(content: bytes, analysis: dict):
    # Note: Figma files are typically accessed via API, not as local files.
    # This is a simplified version assuming we have the JSON content of a Figma file.
    figma_data = json.loads(content)

    analysis['slide_count'] = len(figma_data.get('pages', []))
    analysis['image_count'] = 0
    analysis['bullet_count'] = 0
    analysis['fonts_used'] = set()

    def process_node(node):
        if node['type'] == 'TEXT':
            if node.get('characters', '').strip().startswith(('•', '-', '*')):
                analysis['bullet_count'] += 1
            if 'style' in node and 'fontFamily' in node['style']:
                analysis['fonts_used'].add(node['style']['fontFamily'])
        elif node['type'] == 'IMAGE':
            analysis['image_count'] += 1

        for child in node.get('children', []):
            process_node(child)

    for page in figma_data.get('pages', []):
        process_node(page)

    analysis['fonts_used'] = list(analysis['fonts_used'])


def analyze_readme_md(content: bytes, analysis: dict):
    markdown_content = content.decode('utf-8')

    # Count headings as "slides"
    analysis['slide_count'] = len(
        re.findall(r'^#{1,6}\s', markdown_content, re.MULTILINE))

    # Count bullet points
    analysis['bullet_count'] = len(
        re.findall(r'^\s*[-*+]\s', markdown_content, re.MULTILINE))

    # Count images
    analysis['image_count'] = len(
        re.findall(r'!\[.*?\]\(.*?\)', markdown_content))

    # Markdown doesn't typically include font information
    analysis['fonts_used'] = []

    # Markdown doesn't typically include audio/video, but you could count links to common formats
    analysis['audio_count'] = len(
        re.findall(r'\[.*?\]\(.*?\.(mp3|wav|ogg)\)', markdown_content,
                   re.IGNORECASE))
    analysis['video_count'] = len(
        re.findall(r'\[.*?\]\(.*?\.(mp4|avi|mov)\)', markdown_content,
                   re.IGNORECASE))
